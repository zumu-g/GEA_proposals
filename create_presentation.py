#!/usr/bin/env python3
"""
Premium Marketing Proposal Presentation Generator
171 Greaves Rd, Narre Warren South

STYLE: Luxury Designer Minimalist
Premium color palette + High-impact typography
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# PREMIUM LUXURY COLOR PALETTE
# ============================================================
# Inspired by high-end real estate and luxury brands
# Rich, sophisticated, warm undertones

COLORS = {
    # PRIMARY - Deep sophisticated blacks
    'black': RGBColor(18, 18, 18),           # True luxury black
    'graphite': RGBColor(38, 38, 42),        # Warm graphite (main dark bg)

    # NEUTRALS - Warm sophisticated tones
    'stone': RGBColor(245, 243, 240),        # Warm white/cream (light bg)
    'pearl': RGBColor(252, 251, 249),        # Soft pearl white
    'white': RGBColor(255, 255, 255),        # Pure white for text

    # ACCENT - Luxury metallics
    'champagne': RGBColor(212, 190, 156),    # Champagne gold (primary accent)
    'brass': RGBColor(181, 166, 134),        # Muted brass
    'copper': RGBColor(176, 141, 107),       # Warm copper

    # SECONDARY - Rich earth tones
    'espresso': RGBColor(48, 42, 38),        # Deep espresso brown
    'slate': RGBColor(72, 72, 78),           # Cool slate
    'olive': RGBColor(85, 89, 72),           # Sophisticated olive
    'forest': RGBColor(42, 54, 46),          # Deep forest

    # TEXT
    'text_dark': RGBColor(38, 38, 42),       # Body text on light
    'text_light': RGBColor(245, 243, 240),   # Body text on dark
    'text_muted': RGBColor(140, 138, 135),   # Muted/secondary text
}

# ============================================================
# TYPOGRAPHY
# ============================================================
# Headlines: Didot - elegant, high-fashion serif (like Vogue)
# Subheads: Avenir - clean, sophisticated sans-serif
# Body: Avenir Light - readable, modern

FONTS = {
    'headline': 'Didot',           # Elegant editorial serif
    'subhead': 'Avenir Next',      # Clean modern sans
    'body': 'Avenir Next',         # Readable body
    'accent': 'Avenir Next',       # Labels and captions
}


def create_presentation():
    """Create the full marketing proposal presentation."""
    prs = Presentation()

    # Large widescreen format
    prs.slide_width = Inches(26.67)
    prs.slide_height = Inches(15)

    # Build presentation with rhythm
    add_title_slide(prs)
    add_divider(prs, "we understand")
    add_understanding_slide(prs)
    add_empathy_slide(prs)
    add_divider(prs, "the property")
    add_property_hero_slide(prs)
    add_property_features_slide(prs)
    add_comparable_sales_slide(prs)
    add_divider(prs, "the problem")
    add_problem_slide(prs)
    add_what_you_lose_slide(prs)
    add_divider(prs, "our approach")
    add_approach_slide(prs)
    add_results_slide(prs)
    add_divider(prs, "the campaign")
    add_campaign_overview_slide(prs)
    add_stage1_slide(prs)
    add_stage2_slide(prs)
    add_inspections_slide(prs)
    add_divider(prs, "the investment")
    add_investment_slide(prs)
    add_closing_slide(prs)

    return prs


def add_background(slide, color):
    """Add solid background color to slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, left, top, width, height, text,
             size=24, color=COLORS['text_dark'], bold=False,
             italic=False, align=PP_ALIGN.LEFT, font=FONTS['body']):
    """Add a text box with specified formatting."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    return box


def rect(slide, left, top, width, height, color):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def line(slide, left, top, width, color, thickness=0.05):
    """Add a thin horizontal line."""
    return rect(slide, left, top, width, thickness, color)


# ============================================================
# SLIDE: TITLE
# ============================================================
def add_title_slide(prs):
    """Premium title slide - elegant and bold."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['graphite'])

    # Thin champagne accent line at top
    line(slide, 0, 0, 26.67, COLORS['champagne'], 0.08)

    # Small label top-left - refined
    text_box(slide, 1.5, 2, 10, 0.5,
             "MARKETING PROPOSAL",
             size=14, color=COLORS['brass'], font=FONTS['accent'], bold=False)

    # Main address - large, elegant serif
    text_box(slide, 1.5, 3.5, 22, 2.5,
             "171 Greaves Road",
             size=130, color=COLORS['white'], bold=False, font=FONTS['headline'])

    # Suburb - champagne accent
    text_box(slide, 1.5, 6.5, 15, 1,
             "Narre Warren South",
             size=42, color=COLORS['champagne'], font=FONTS['headline'], italic=True)

    # Thin line separator
    line(slide, 1.5, 8.5, 8, COLORS['champagne'], 0.03)

    # Property highlights - refined
    text_box(slide, 1.5, 9.5, 20, 1,
             "Designer Residence  ·  1 Acre  ·  37 Squares  ·  Resort Pool  ·  216sqm Shed",
             size=18, color=COLORS['text_muted'], font=FONTS['body'])

    # Bottom bar
    rect(slide, 0, 12, 26.67, 3, COLORS['espresso'])

    # Agency placeholder
    text_box(slide, 1.5, 13, 10, 0.5,
             "[YOUR AGENCY]",
             size=14, color=COLORS['brass'], font=FONTS['accent'])


# ============================================================
# SLIDE: DIVIDER
# ============================================================
def add_divider(prs, text):
    """Section divider - elegant and impactful."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['black'])

    # Champagne accent lines
    line(slide, 0, 0, 26.67, COLORS['champagne'], 0.06)
    line(slide, 0, 14.94, 26.67, COLORS['champagne'], 0.06)

    # Main text - elegant serif, slightly smaller for refinement
    text_box(slide, 1.5, 5.5, 24, 3,
             text,
             size=110, color=COLORS['white'], bold=False, font=FONTS['headline'])


# ============================================================
# SLIDE: UNDERSTANDING (Accusation Audit)
# ============================================================
def add_understanding_slide(prs):
    """We know what you're thinking - split layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['stone'])

    # Headline - elegant serif
    text_box(slide, 1.5, 2, 12, 3,
             "We know what\nyou're thinking",
             size=68, color=COLORS['text_dark'], bold=False, font=FONTS['headline'])

    # Thin accent line under headline
    line(slide, 1.5, 5.5, 6, COLORS['champagne'], 0.04)

    # Content - right side, refined body text
    lines = [
        "You've invested years creating something exceptional.",
        "",
        "You're wondering if any agent truly understands",
        "the value of what you've built.",
        "",
        "You're concerned that standard marketing",
        "won't capture what makes this special.",
    ]

    y = 3.5
    for ln in lines:
        if ln:
            text_box(slide, 14, y, 11, 0.8,
                     ln,
                     size=24, color=COLORS['text_dark'], font=FONTS['body'])
        y += 0.85


# ============================================================
# SLIDE: EMPATHY (Labeling)
# ============================================================
def add_empathy_slide(prs):
    """It seems like... It sounds like..."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['graphite'])

    line(slide, 0, 0, 26.67, COLORS['champagne'], 0.06)

    # Headline
    text_box(slide, 1.5, 1.5, 20, 1.5,
             "This isn't just a property",
             size=64, color=COLORS['white'], bold=False, font=FONTS['headline'])

    # The labeling statements - elegant layout
    statements = [
        ("It seems like", "you've created more than a home—a private sanctuary"),
        ("It sounds like", "every detail was intentional, every finish considered"),
        ("It feels like", "this deserves a buyer who truly appreciates the vision"),
    ]

    y = 5
    for prefix, content in statements:
        text_box(slide, 1.5, y, 5, 0.8,
                 prefix,
                 size=22, color=COLORS['champagne'], italic=True, font=FONTS['headline'])
        text_box(slide, 7, y, 18, 0.8,
                 content,
                 size=22, color=COLORS['white'], font=FONTS['body'])
        y += 2.2

    # Bottom bar
    rect(slide, 0, 12.5, 26.67, 2.5, COLORS['forest'])
    text_box(slide, 1.5, 13.3, 24, 1,
             "We're here to ensure the market sees what you see.",
             size=24, color=COLORS['text_light'], italic=True, font=FONTS['subhead'])


# ============================================================
# SLIDE: PROPERTY HERO
# ============================================================
def add_property_hero_slide(prs):
    """Property hero - image placeholder with elegant typography."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['stone'])

    # Image placeholder - left side
    rect(slide, 0, 0, 14, 15, COLORS['slate'])
    text_box(slide, 4.5, 7, 5, 1,
             "[HERO IMAGE]",
             size=20, color=COLORS['text_muted'], align=PP_ALIGN.CENTER, font=FONTS['accent'])

    # Content - right side
    text_box(slide, 15.5, 2, 10, 1,
             "The Property",
             size=56, color=COLORS['text_dark'], bold=False, font=FONTS['headline'])

    line(slide, 15.5, 3.8, 4, COLORS['champagne'], 0.04)

    text_box(slide, 15.5, 5, 10, 4,
             "A designer custom-built residence\non a full acre of private gardens.\n\nThe space, privacy, and quality\nthat discerning buyers seek\nbut rarely find.",
             size=22, color=COLORS['text_dark'], font=FONTS['body'])

    # Key stats - large elegant numbers
    text_box(slide, 15.5, 10, 3, 1.5,
             "37",
             size=72, color=COLORS['copper'], bold=False, font=FONTS['headline'])
    text_box(slide, 15.5, 11.5, 3, 0.5,
             "squares",
             size=16, color=COLORS['text_muted'], font=FONTS['accent'])

    text_box(slide, 20, 10, 3, 1.5,
             "1",
             size=72, color=COLORS['copper'], bold=False, font=FONTS['headline'])
    text_box(slide, 20, 11.5, 3, 0.5,
             "acre",
             size=16, color=COLORS['text_muted'], font=FONTS['accent'])


# ============================================================
# SLIDE: PROPERTY FEATURES
# ============================================================
def add_property_features_slide(prs):
    """Property features - clean elegant grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['graphite'])

    line(slide, 0, 0, 26.67, COLORS['champagne'], 0.06)

    # Headline
    text_box(slide, 1.5, 1.5, 10, 1,
             "The Details",
             size=56, color=COLORS['white'], bold=False, font=FONTS['headline'])

    # Two columns of features
    left_features = [
        ("37 Squares", "Designer living spaces"),
        ("Formal Lounge", "With fireplace"),
        ("Teenagers Retreat", "Zoned living"),
        ("900mm Cooktop", "Premium kitchen"),
        ("7.5 Star Energy", "Double glazed throughout"),
    ]

    right_features = [
        ("1 Acre", "Private grounds"),
        ("Solar Heated Pool", "With cabana"),
        ("216sqm Shed", "3-phase power"),
        ("Built-in BBQ", "Alfresco kitchen"),
        ("Serene Gardens", "Established landscaping"),
    ]

    y = 4.5
    for (title, desc) in left_features:
        text_box(slide, 1.5, y, 8, 0.6,
                 title,
                 size=24, color=COLORS['white'], bold=False, font=FONTS['subhead'])
        text_box(slide, 1.5, y + 0.65, 8, 0.5,
                 desc,
                 size=16, color=COLORS['text_muted'], font=FONTS['body'])
        y += 1.7

    y = 4.5
    for (title, desc) in right_features:
        text_box(slide, 14, y, 8, 0.6,
                 title,
                 size=24, color=COLORS['white'], bold=False, font=FONTS['subhead'])
        text_box(slide, 14, y + 0.65, 8, 0.5,
                 desc,
                 size=16, color=COLORS['text_muted'], font=FONTS['body'])
        y += 1.7

    # Elegant divider line
    rect(slide, 12, 4.5, 0.03, 8, COLORS['champagne'])


# ============================================================
# SLIDE: COMPARABLE SALES
# ============================================================
def add_comparable_sales_slide(prs):
    """Recent nearby sales - elegant data presentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['stone'])

    # Headline
    text_box(slide, 1.5, 1.5, 15, 1,
             "Market Evidence",
             size=56, color=COLORS['text_dark'], bold=False, font=FONTS['headline'])

    text_box(slide, 1.5, 3.2, 15, 0.8,
             "Recent sales on Greaves Road",
             size=22, color=COLORS['text_muted'], font=FONTS['body'])

    line(slide, 1.5, 4.2, 5, COLORS['champagne'], 0.04)

    # Sales cards - elegant dark cards
    sales = [
        ("165 Greaves Road", "$1,890,000", "4 bed  ·  2 bath  ·  10 car"),
        ("173 Greaves Road", "$1,780,000", "4 bed  ·  2 bath  ·  7 car"),
        ("197-205 Greaves Road", "$1,700,000", "4 bed  ·  2 bath  ·  3 car"),
    ]

    x = 1.5
    for (address, price, features) in sales:
        rect(slide, x, 5.2, 7.5, 4.5, COLORS['graphite'])
        text_box(slide, x + 0.6, 5.8, 6.3, 0.8,
                 address,
                 size=18, color=COLORS['text_light'], font=FONTS['accent'])
        text_box(slide, x + 0.6, 6.8, 6.3, 1.2,
                 price,
                 size=48, color=COLORS['champagne'], bold=False, font=FONTS['headline'])
        text_box(slide, x + 0.6, 8.5, 6.3, 0.6,
                 features,
                 size=14, color=COLORS['text_muted'], font=FONTS['body'])
        x += 8.2

    # Suburb stats at bottom
    rect(slide, 0, 11, 26.67, 4, COLORS['espresso'])

    stats = [
        ("$820K", "Median Price"),
        ("21", "Days on Market"),
        ("400+", "Sold (12 months)"),
        ("+1.8%", "Annual Growth"),
    ]

    x = 2.5
    for (value, label) in stats:
        text_box(slide, x, 11.8, 5, 1.2,
                 value,
                 size=44, color=COLORS['champagne'], bold=False, font=FONTS['headline'])
        text_box(slide, x, 13.3, 5, 0.5,
                 label,
                 size=14, color=COLORS['text_muted'], font=FONTS['accent'])
        x += 6


# ============================================================
# SLIDE: THE PROBLEM
# ============================================================
def add_problem_slide(prs):
    """What goes wrong - elegant split layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['stone'])

    # Headline
    text_box(slide, 1.5, 2, 12, 3,
             "Ordinary marketing\ngets ordinary results",
             size=58, color=COLORS['text_dark'], bold=False, font=FONTS['headline'])

    line(slide, 1.5, 5.8, 5, COLORS['copper'], 0.04)

    # Problems list - right side
    problems = [
        "Generic photography misses the details",
        "Template marketing attracts wrong buyers",
        "Open homes become a parade of lookers",
        "Days on market erode perceived value",
    ]

    y = 4
    for problem in problems:
        text_box(slide, 14, y, 11, 1,
                 "—  " + problem,
                 size=20, color=COLORS['text_dark'], font=FONTS['body'])
        y += 1.6

    # Callout box
    rect(slide, 14, 10.5, 11, 3, COLORS['graphite'])
    text_box(slide, 14.6, 11.3, 10, 2,
             "You only get one chance to launch.",
             size=24, color=COLORS['champagne'], italic=True, font=FONTS['subhead'])


# ============================================================
# SLIDE: WHAT YOU LOSE
# ============================================================
def add_what_you_lose_slide(prs):
    """Loss aversion - bold elegant statements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['graphite'])

    line(slide, 0, 0, 26.67, COLORS['champagne'], 0.06)

    # Headline
    text_box(slide, 1.5, 1.5, 15, 1,
             "What Gets Lost",
             size=64, color=COLORS['white'], bold=False, font=FONTS['headline'])

    # Losses with elegant accent bars
    losses = [
        "The craftsmanship becomes invisible",
        "The lifestyle story isn't told",
        "Premium buyers scroll past",
        "You negotiate with bargain hunters",
    ]

    y = 5
    for loss in losses:
        rect(slide, 1.5, y + 0.15, 0.15, 0.7, COLORS['champagne'])
        text_box(slide, 2.2, y, 22, 1,
                 loss,
                 size=32, color=COLORS['white'], font=FONTS['body'])
        y += 1.7

    # Bottom quote
    rect(slide, 0, 12, 26.67, 3, COLORS['espresso'])
    text_box(slide, 1.5, 12.8, 24, 1.5,
             '"Your property is only as valuable as buyers perceive it to be."',
             size=26, color=COLORS['champagne'], italic=True, font=FONTS['subhead'])


# ============================================================
# SLIDE: OUR APPROACH
# ============================================================
def add_approach_slide(prs):
    """How we do things differently - elegant numbered list."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['stone'])

    # Headline
    text_box(slide, 1.5, 1.5, 15, 1,
             "How We Work",
             size=64, color=COLORS['text_dark'], bold=False, font=FONTS['headline'])

    line(slide, 1.5, 3.3, 4, COLORS['champagne'], 0.04)

    # Approach points
    points = [
        ("01", "Custom Creative", "Bespoke marketing designed for your property"),
        ("02", "Strategic Exposure", "Targeting the right buyers, not all buyers"),
        ("03", "Exclusivity", "Private showings that signal value"),
        ("04", "Premium Presentation", "Professional staging, twilight photography, cinematic video"),
        ("05", "Expert Negotiation", "Creating competition to maximise your outcome"),
    ]

    y = 4.5
    for (num, title, desc) in points:
        text_box(slide, 1.5, y, 1.5, 1,
                 num,
                 size=36, color=COLORS['copper'], bold=False, font=FONTS['headline'])
        text_box(slide, 4, y - 0.1, 10, 0.8,
                 title,
                 size=24, color=COLORS['text_dark'], bold=False, font=FONTS['subhead'])
        text_box(slide, 4, y + 0.7, 18, 0.6,
                 desc,
                 size=16, color=COLORS['text_muted'], font=FONTS['body'])
        y += 1.9


# ============================================================
# SLIDE: RESULTS
# ============================================================
def add_results_slide(prs):
    """Track record - elegant layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['graphite'])

    line(slide, 0, 0, 26.67, COLORS['champagne'], 0.06)

    # Headline
    text_box(slide, 1.5, 1.5, 15, 1,
             "Proven Results",
             size=64, color=COLORS['white'], bold=False, font=FONTS['headline'])

    # Case study box
    rect(slide, 1.5, 4, 13, 7, COLORS['espresso'])
    text_box(slide, 2.1, 4.6, 12, 0.6,
             "CASE STUDY",
             size=12, color=COLORS['champagne'], font=FONTS['accent'])
    text_box(slide, 2.1, 5.8, 12, 5,
             "A comparable acreage property was\nmeticulously marketed with professional\nstaging, premium visuals, and strategic\nVIP outreach.\n\nResult: Sale price exceeding\nvendor expectations.",
             size=24, color=COLORS['white'], font=FONTS['body'])

    # Placeholder for stats
    rect(slide, 16, 4, 9, 7, COLORS['stone'])
    text_box(slide, 16.6, 4.6, 8, 0.6,
             "YOUR RESULTS",
             size=12, color=COLORS['espresso'], font=FONTS['accent'])
    text_box(slide, 16.6, 6.5, 8, 4,
             "[Insert your specific\nresults and testimonials]",
             size=18, color=COLORS['text_muted'], font=FONTS['body'])


# ============================================================
# SLIDE: CAMPAIGN OVERVIEW
# ============================================================
def add_campaign_overview_slide(prs):
    """Two-stage campaign overview - elegant cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['stone'])

    # Headline
    text_box(slide, 1.5, 1.5, 15, 1,
             "Two Stages",
             size=64, color=COLORS['text_dark'], bold=False, font=FONTS['headline'])

    line(slide, 1.5, 3.3, 3, COLORS['champagne'], 0.04)

    # Stage 1
    rect(slide, 1.5, 4.5, 11, 8.5, COLORS['graphite'])
    text_box(slide, 2.1, 5.2, 10, 0.6,
             "STAGE 1",
             size=12, color=COLORS['champagne'], font=FONTS['accent'])
    text_box(slide, 2.1, 6.2, 10, 1.5,
             "VIP Campaign",
             size=44, color=COLORS['white'], bold=False, font=FONTS['headline'])
    text_box(slide, 2.1, 8.2, 10, 0.6,
             "2 weeks  ·  Off-market",
             size=14, color=COLORS['text_muted'], font=FONTS['accent'])
    text_box(slide, 2.1, 9.5, 10, 2.5,
             "Exclusive preview for qualified buyers.\nCreate urgency and competition before\nthe public even knows you're selling.",
             size=18, color=COLORS['text_light'], font=FONTS['body'])

    # Elegant arrow
    text_box(slide, 13, 8, 1, 1,
             "→",
             size=48, color=COLORS['champagne'], font=FONTS['body'])

    # Stage 2
    rect(slide, 14.5, 4.5, 11, 8.5, COLORS['espresso'])
    text_box(slide, 15.1, 5.2, 10, 0.6,
             "STAGE 2",
             size=12, color=COLORS['champagne'], font=FONTS['accent'])
    text_box(slide, 15.1, 6.2, 10, 1.5,
             "Public Launch",
             size=44, color=COLORS['white'], bold=False, font=FONTS['headline'])
    text_box(slide, 15.1, 8.2, 10, 0.6,
             "Ongoing  ·  Maximum exposure",
             size=14, color=COLORS['text_muted'], font=FONTS['accent'])
    text_box(slide, 15.1, 9.5, 10, 2.5,
             "Full premium campaign with targeted\ndigital advertising and direct outreach\nto qualified buyers.",
             size=18, color=COLORS['text_light'], font=FONTS['body'])


# ============================================================
# SLIDE: STAGE 1 DETAIL
# ============================================================
def add_stage1_slide(prs):
    """VIP Campaign detail."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['graphite'])

    line(slide, 0, 0, 26.67, COLORS['champagne'], 0.06)

    # Headline
    text_box(slide, 1.5, 1.5, 15, 1,
             "Stage 1",
             size=44, color=COLORS['white'], bold=False, font=FONTS['headline'])

    text_box(slide, 1.5, 2.8, 15, 1,
             "Discreet VIP Campaign",
             size=32, color=COLORS['champagne'], font=FONTS['subhead'])

    # Two columns
    rect(slide, 1.5, 4.5, 11, 8.5, COLORS['espresso'])
    text_box(slide, 2.1, 5.1, 10, 0.5,
             "WHAT WE DO",
             size=12, color=COLORS['champagne'], font=FONTS['accent'])

    tactics = [
        "Personal outreach to VIP buyers",
        "Private evening showings",
        "Gourmet catering experience",
        "Buyer's agent network alerted",
        "Teaser social campaign",
    ]
    y = 6.2
    for tactic in tactics:
        text_box(slide, 2.1, y, 10, 0.7,
                 "·  " + tactic,
                 size=18, color=COLORS['white'], font=FONTS['body'])
        y += 1.1

    rect(slide, 14, 4.5, 11, 8.5, COLORS['stone'])
    text_box(slide, 14.6, 5.1, 10, 0.5,
             "THE OUTCOME",
             size=12, color=COLORS['espresso'], font=FONTS['accent'])

    outcomes = [
        "Urgency before public launch",
        "Competitive environment created",
        "Serious buyers identified",
        "Offers potentially secured early",
        "Property positioned as exclusive",
    ]
    y = 6.2
    for outcome in outcomes:
        text_box(slide, 14.6, y, 10, 0.7,
                 "✓  " + outcome,
                 size=18, color=COLORS['text_dark'], font=FONTS['body'])
        y += 1.1


# ============================================================
# SLIDE: STAGE 2 DETAIL
# ============================================================
def add_stage2_slide(prs):
    """Public campaign detail."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['stone'])

    # Left accent bar
    rect(slide, 0, 0, 0.8, 15, COLORS['espresso'])

    # Headline
    text_box(slide, 2, 1.5, 15, 1,
             "Stage 2",
             size=44, color=COLORS['text_dark'], bold=False, font=FONTS['headline'])

    text_box(slide, 2, 2.8, 15, 1,
             "Premium Marketing Campaign",
             size=28, color=COLORS['text_muted'], font=FONTS['subhead'])

    line(slide, 2, 4, 4, COLORS['champagne'], 0.04)

    # Visual marketing
    text_box(slide, 2, 5, 10, 0.5,
             "VISUAL MARKETING",
             size=12, color=COLORS['espresso'], font=FONTS['accent'])

    visual = ["Professional photography", "Cinematic drone footage",
              "Video tour", "Brochure with floorplans", "Premium signboard"]
    y = 5.8
    for item in visual:
        text_box(slide, 2, y, 10, 0.6,
                 "·  " + item,
                 size=17, color=COLORS['text_dark'], font=FONTS['body'])
        y += 0.8

    # Digital campaign
    text_box(slide, 14, 5, 10, 0.5,
             "DIGITAL CAMPAIGN",
             size=12, color=COLORS['espresso'], font=FONTS['accent'])

    digital = ["realestate.com.au Premiere+", "Domain Platinum",
               "Facebook & Instagram targeting", "Google remarketing", "YouTube pre-roll"]
    y = 5.8
    for item in digital:
        text_box(slide, 14, y, 10, 0.6,
                 "·  " + item,
                 size=17, color=COLORS['text_dark'], font=FONTS['body'])
        y += 0.8

    # Bottom bar
    rect(slide, 0, 11.5, 26.67, 3.5, COLORS['graphite'])
    text_box(slide, 2, 12.5, 23, 1,
             "Maximum exposure to qualified buyers.",
             size=24, color=COLORS['champagne'], italic=True, font=FONTS['subhead'])


# ============================================================
# SLIDE: INSPECTIONS
# ============================================================
def add_inspections_slide(prs):
    """Private inspection protocol - elegant grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['graphite'])

    line(slide, 0, 0, 26.67, COLORS['champagne'], 0.06)

    # Headline
    text_box(slide, 1.5, 1.5, 20, 1,
             "Private Showings",
             size=64, color=COLORS['white'], bold=False, font=FONTS['headline'])

    text_box(slide, 1.5, 3.3, 20, 1,
             "By appointment only",
             size=26, color=COLORS['champagne'], italic=True, font=FONTS['subhead'])

    # Benefits grid
    benefits = [
        ("Pre-qualified", "Genuine buyers only"),
        ("Curated", "Premium experience"),
        ("Presentation", "Property at its best"),
        ("Exclusivity", "Buyers compete"),
    ]

    x = 1.5
    for (title, desc) in benefits:
        rect(slide, x, 5.5, 5.8, 4.5, COLORS['espresso'])
        text_box(slide, x + 0.5, 6.5, 4.8, 0.8,
                 title,
                 size=26, color=COLORS['champagne'], bold=False, font=FONTS['subhead'])
        text_box(slide, x + 0.5, 7.8, 4.8, 1,
                 desc,
                 size=18, color=COLORS['text_light'], font=FONTS['body'])
        x += 6.2


# ============================================================
# SLIDE: INVESTMENT
# ============================================================
def add_investment_slide(prs):
    """Marketing investment breakdown - elegant and clear."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['stone'])

    # Headline
    text_box(slide, 1.5, 1.5, 15, 1,
             "The Investment",
             size=64, color=COLORS['text_dark'], bold=False, font=FONTS['headline'])

    line(slide, 1.5, 3.3, 4, COLORS['champagne'], 0.04)

    # Investment items
    items = [
        ("Photography & Video", "$X,XXX"),
        ("Staging (if required)", "$X,XXX"),
        ("Signboard & Brochures", "$XXX"),
        ("realestate.com.au Premiere+", "$X,XXX"),
        ("Domain Platinum", "$XXX"),
        ("Social Media Advertising", "$X,XXX"),
        ("VIP Event Catering", "$XXX"),
    ]

    y = 4.5
    for (item, cost) in items:
        text_box(slide, 1.5, y, 9, 0.7,
                 item,
                 size=20, color=COLORS['text_dark'], font=FONTS['body'])
        text_box(slide, 10.5, y, 3, 0.7,
                 cost,
                 size=20, color=COLORS['text_dark'], align=PP_ALIGN.RIGHT, font=FONTS['body'])
        y += 1.0

    # Total line
    line(slide, 1.5, y + 0.2, 12, COLORS['champagne'], 0.03)
    text_box(slide, 1.5, y + 0.5, 9, 0.8,
             "Total",
             size=24, color=COLORS['text_dark'], bold=False, font=FONTS['subhead'])
    text_box(slide, 10.5, y + 0.5, 3, 0.8,
             "$XX,XXX",
             size=24, color=COLORS['text_dark'], bold=False, align=PP_ALIGN.RIGHT, font=FONTS['subhead'])

    # Commission box
    rect(slide, 16, 4.5, 9, 6, COLORS['graphite'])
    text_box(slide, 16.6, 5.1, 8, 0.5,
             "COMMISSION",
             size=12, color=COLORS['champagne'], font=FONTS['accent'])
    text_box(slide, 16.6, 6.5, 8, 3,
             "X.X% + GST\n\nPayable on\nsuccessful settlement",
             size=24, color=COLORS['white'], font=FONTS['body'])


# ============================================================
# SLIDE: CLOSING
# ============================================================
def add_closing_slide(prs):
    """The close - elegant and confident."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS['black'])

    line(slide, 0, 0, 26.67, COLORS['champagne'], 0.06)
    line(slide, 0, 14.94, 26.67, COLORS['champagne'], 0.06)

    # Main headline - elegant serif
    text_box(slide, 1.5, 3.5, 24, 2.5,
             "The Next Step",
             size=100, color=COLORS['white'], bold=False, font=FONTS['headline'])

    # Calibrated question
    text_box(slide, 1.5, 7.5, 24, 1.5,
             '"How would you like to proceed?"',
             size=36, color=COLORS['champagne'], italic=True, font=FONTS['headline'])

    # Next steps box
    rect(slide, 8, 10, 11, 3.5, COLORS['espresso'])
    text_box(slide, 8.6, 10.5, 10, 0.5,
             "WHEN YOU'RE READY",
             size=12, color=COLORS['champagne'], font=FONTS['accent'])
    text_box(slide, 8.6, 11.5, 10, 2,
             "1. Sign agreement    2. Schedule photography\n3. Confirm marketing    4. Launch campaign",
             size=18, color=COLORS['white'], font=FONTS['body'])

    # Contact
    text_box(slide, 1.5, 13.5, 6, 0.5,
             "[Your Name]  ·  [Phone]  ·  [Email]",
             size=14, color=COLORS['text_muted'], font=FONTS['accent'])


# ============================================================
# MAIN
# ============================================================
if __name__ == "__main__":
    output_dir = "/Users/stuartgrant_mbp13/Library/Mobile Documents/com~apple~CloudDocs/GEA_proposals"
    output_file = os.path.join(output_dir, "171_Greaves_Rd_Marketing_Proposal.pptx")

    print("Creating presentation with premium designer style...")
    print("\nColor Palette:")
    print("  · Black (#121212) - Luxury black")
    print("  · Graphite (#26262A) - Primary dark")
    print("  · Stone (#F5F3F0) - Warm light")
    print("  · Champagne (#D4BE9C) - Primary accent")
    print("  · Copper (#B08D6B) - Secondary accent")
    print("  · Espresso (#302A26) - Deep brown")
    print("\nTypography:")
    print("  · Headlines: Didot (elegant editorial serif)")
    print("  · Subheads: Avenir Next (clean modern)")
    print("  · Body: Avenir Next (readable)")

    prs = create_presentation()

    print(f"\nSaving to: {output_file}")
    prs.save(output_file)

    print("\nDone!")
    print(f"Slide count: {len(prs.slides)}")
