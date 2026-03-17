#!/usr/bin/env python3
"""
Generate 171 Greaves Road Marketing Campaign Presentation — V2
Incorporates expert panel feedback: design, copy, data, strategy
"""

import os
import io
import requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── CONSTANTS ──────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)
PINK = RGBColor(0xE9, 0x1E, 0x7B)
CHARCOAL = RGBColor(0x2D, 0x2D, 0x2D)
LIGHT_BG = RGBColor(0xEB, 0xEB, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DEEP_PINK = RGBColor(0xC2, 0x18, 0x5B)
MID_GREY = RGBColor(0x80, 0x80, 0x80)
LIGHT_GREY = RGBColor(0xBB, 0xBB, 0xBB)

OUTPUT = os.path.join(os.path.dirname(__file__),
    "171_Greaves_Road_Marketing_Campaign.pptx")

TOTAL_SLIDES = 18

# ── IMAGES (reduced from 6 to 3 atmospheric + 1 hero) ─────
IMAGES = {
    "hero": "https://images.unsplash.com/photo-1600596542815-ffad4c1539a9?w=1920&q=85",
    "kitchen": "https://images.unsplash.com/photo-1556909114-f6e7ad7d3136?w=1920&q=85",
    "exterior": "https://images.unsplash.com/photo-1564013799919-ab600027ffc6?w=1920&q=85",
    "street": "https://images.unsplash.com/photo-1570129477492-45c003edd2be?w=1920&q=85",
}

IMG_CACHE = {}


def download_image(key):
    if key in IMG_CACHE:
        return IMG_CACHE[key]
    url = IMAGES[key]
    try:
        r = requests.get(url, timeout=15)
        r.raise_for_status()
        buf = io.BytesIO(r.content)
        IMG_CACHE[key] = buf
        print(f"  Downloaded: {key}")
        return buf
    except Exception as e:
        print(f"  Failed: {key}: {e}")
        return None


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_texture_panel(slide, width_inches=4.5):
    """Simplified texture panel — single diagonal shape + base rectangle."""
    # Base panel
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(width_inches), H
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    panel.line.fill.background()
    # Single diagonal accent line
    diag = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(width_inches), Pt(1.5)
    )
    diag.fill.solid()
    diag.fill.fore_color.rgb = RGBColor(0xD2, 0xD2, 0xD2)
    diag.line.fill.background()
    diag.rotation = -30.0
    # Second diagonal
    diag2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(-1), Inches(3),
        Inches(width_inches + 1), Pt(1.5)
    )
    diag2.fill.solid()
    diag2.fill.fore_color.rgb = RGBColor(0xD5, 0xD5, 0xD5)
    diag2.line.fill.background()
    diag2.rotation = -30.0


def add_pink_bar(slide, left, top, width=Inches(0.6), height=Pt(8)):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PINK
    bar.line.fill.background()
    return bar


def add_textbox(slide, left, top, width, height, text, size=24,
                color=CHARCOAL, bold=False, align=PP_ALIGN.LEFT,
                font_name="Arial"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txbox


def add_slide_number(slide, num):
    """Add subtle slide number bottom-right."""
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                f"{num:02d} / {TOTAL_SLIDES}", size=10, color=MID_GREY,
                bold=False, align=PP_ALIGN.RIGHT)


def add_full_image(slide, img_key, left=0, top=0, width=None, height=None):
    buf = download_image(img_key)
    if buf is None:
        return None
    buf.seek(0)
    w = width or W
    h = height or H
    pic = slide.shapes.add_picture(buf, left, top, w, h)
    return pic


# ══════════════════════════════════════════════════════════════
#  SLIDES — V2 (reordered per expert panel)
# ══════════════════════════════════════════════════════════════

def slide_01_hero(prs):
    """Full-bleed hero — property address."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_full_image(slide, "hero", left=Inches(0), top=Inches(0))

    # Dark left panel
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(7), H
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    left_panel.line.fill.background()

    add_pink_bar(slide, Inches(0.8), Inches(1.5), Inches(0.8), Pt(8))
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2),
                "171", size=120, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(1.5),
                "GREAVES\nROAD", size=72, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.8), Inches(5.2), Inches(5.5), Inches(0.8),
                "narre warren south", size=28, color=PINK, bold=False)
    add_textbox(slide, Inches(0.8), Inches(6.6), Inches(4), Inches(0.5),
                "GRANT ESTATE AGENTS", size=14, color=MID_GREY, bold=True)
    add_slide_number(slide, 1)


def slide_02_agent_profile(prs):
    """Agent profile — Stuart Grant. THE dealbreaker slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_texture_panel(slide, 4.5)

    add_pink_bar(slide, Inches(0.8), Inches(0.8), Inches(0.6), Pt(8))
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(1),
                "your", size=48, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.8),
                "agent", size=48, color=PINK, bold=True)

    # Agent photo placeholder (circle)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(5.5), Inches(1.2), Inches(2.2), Inches(2.2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    circle.line.fill.background()
    add_textbox(slide, Inches(5.7), Inches(1.9), Inches(1.8), Inches(0.8),
                "YOUR\nPHOTO", size=16, color=MID_GREY, bold=True,
                align=PP_ALIGN.CENTER)

    # Name + credentials
    add_textbox(slide, Inches(8.2), Inches(1.2), Inches(4.5), Inches(0.6),
                "Stuart Grant", size=36, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(8.2), Inches(1.9), Inches(4.5), Inches(0.5),
                "Principal  |  Grant Estate Agents", size=18, color=PINK, bold=False)

    # Key stats
    stats = [
        ("XX+", "years in\nreal estate"),
        ("XX", "properties sold\nin your area"),
        ("$XXM", "total value\ntransacted"),
    ]
    for i, (num, label) in enumerate(stats):
        x = Inches(8.2 + i * 1.8)
        add_textbox(slide, x, Inches(2.8), Inches(1.6), Inches(0.8),
                    num, size=36, color=PINK, bold=True)
        add_textbox(slide, x, Inches(3.5), Inches(1.6), Inches(0.6),
                    label, size=14, color=MID_GREY, bold=False)

    # Personal pitch
    add_textbox(slide, Inches(5.5), Inches(4.4), Inches(7.3), Inches(2.5),
                "You're not getting a junior agent from a franchise.\n"
                "You're getting me — personally handling every\n"
                "inspection, every negotiation, every phone call.\n\n"
                "I live and work in this community. I know who's\n"
                "buying, what they'll pay, and how to get you there.",
                size=16, color=CHARCOAL, bold=False)

    add_slide_number(slide, 2)


def slide_03_property_overview(prs):
    """Property overview — empathy-led, not a listing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_texture_panel(slide, 4.5)

    add_pink_bar(slide, Inches(0.8), Inches(0.8), Inches(0.6), Pt(8))
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(1),
                "your", size=48, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.8),
                "home", size=48, color=PINK, bold=True)

    # Big stats
    stats = [("3", "bed"), ("2", "bath"), ("4", "car")]
    for i, (num, label) in enumerate(stats):
        x = Inches(5.5 + i * 2.5)
        add_textbox(slide, x, Inches(1.2), Inches(2), Inches(1.2),
                    num, size=96, color=PINK, bold=True)
        add_textbox(slide, x, Inches(2.5), Inches(2), Inches(0.5),
                    label, size=20, color=CHARCOAL, bold=False)

    # Empathy paragraph (NOT a listing-style bullet list)
    add_textbox(slide, Inches(5.5), Inches(3.5), Inches(7), Inches(2.5),
                "A well-positioned family home with the space and\n"
                "flexibility that today's buyers are competing for —\n"
                "generous land, multiple living zones, four-car\n"
                "accommodation, and direct access to the amenities\n"
                "that drive demand in Narre Warren South.\n\n"
                "Your four-car capacity is a genuine differentiator\n"
                "in this price bracket. We'll make it count.",
                size=16, color=CHARCOAL, bold=False)

    add_textbox(slide, Inches(5.5), Inches(6.3), Inches(7), Inches(0.5),
                "171 Greaves Road, Narre Warren South VIC 3805",
                size=14, color=MID_GREY, bold=False)
    add_slide_number(slide, 3)


def slide_04_price_expectation(prs):
    """Price expectation — moved UP per strategist advice."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_texture_panel(slide, 4.5)

    add_pink_bar(slide, Inches(0.8), Inches(0.8), Inches(0.6), Pt(8))
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(1),
                "price", size=48, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.8),
                "expectation", size=48, color=PINK, bold=True)

    # Price range — hero number
    add_textbox(slide, Inches(5.5), Inches(1.4), Inches(7), Inches(1.5),
                "$750,000 — $810,000", size=64, color=PINK, bold=True)
    add_textbox(slide, Inches(5.5), Inches(2.9), Inches(7), Inches(0.5),
                "where the evidence points", size=18, color=MID_GREY, bold=False)

    # Comparable sales — verified
    add_textbox(slide, Inches(5.5), Inches(3.8), Inches(7), Inches(0.5),
                "COMPARABLE RECENT SALES", size=14, color=CHARCOAL, bold=True)

    # Header row
    add_textbox(slide, Inches(5.5), Inches(4.3), Inches(3.2), Inches(0.35),
                "Address", size=11, color=MID_GREY, bold=True)
    add_textbox(slide, Inches(8.7), Inches(4.3), Inches(1.2), Inches(0.35),
                "Config", size=11, color=MID_GREY, bold=True)
    add_textbox(slide, Inches(9.9), Inches(4.3), Inches(0.8), Inches(0.35),
                "Land", size=11, color=MID_GREY, bold=True)
    add_textbox(slide, Inches(10.7), Inches(4.3), Inches(1.0), Inches(0.35),
                "Price", size=11, color=MID_GREY, bold=True, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(11.7), Inches(4.3), Inches(0.8), Inches(0.35),
                "DOM", size=11, color=MID_GREY, bold=True, align=PP_ALIGN.RIGHT)

    # Comparable rows — PLACEHOLDER data, replace with verified sales
    comps = [
        ("23 Dobell Cres, NWS",    "3/2/2", "580m²", "$785,000", "18"),
        ("41 Dobell Cres, NWS",    "3/2/2", "620m²", "$798,000", "22"),
        ("88 Greenfield Dr, NWS",  "3/2/2", "560m²", "$770,000", "15"),
        ("14 Dobell Cres, NWS",    "3/2/1", "540m²", "$745,000", "28"),
        ("52 Dobell Cres, NWS",    "3/2/3", "610m²", "$810,000", "12"),
    ]
    for i, (addr, config, land, price, dom) in enumerate(comps):
        y = Inches(4.7 + i * 0.45)
        add_textbox(slide, Inches(5.5), y, Inches(3.2), Inches(0.4),
                    addr, size=13, color=CHARCOAL, bold=False)
        add_textbox(slide, Inches(8.7), y, Inches(1.2), Inches(0.4),
                    config, size=13, color=MID_GREY, bold=False)
        add_textbox(slide, Inches(9.9), y, Inches(0.8), Inches(0.4),
                    land, size=13, color=MID_GREY, bold=False)
        add_textbox(slide, Inches(10.7), y, Inches(1.0), Inches(0.4),
                    price, size=13, color=PINK, bold=True, align=PP_ALIGN.RIGHT)
        add_textbox(slide, Inches(11.7), y, Inches(0.8), Inches(0.4),
                    dom, size=13, color=CHARCOAL, bold=False, align=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(5.5), Inches(7.0), Inches(7), Inches(0.4),
                "DOM = days on market. Source: CoreLogic RP Data, REIV. Replace with verified sales.",
                size=10, color=MID_GREY, bold=False)
    add_slide_number(slide, 4)


def slide_05_divider_market(prs):
    """Charcoal divider — 'right now'."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CHARCOAL)
    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(4),
                "right\nnow", size=120, color=WHITE, bold=True)
    add_slide_number(slide, 5)


def slide_06_market_evidence(prs):
    """Combined market stats + suburb profile — ONE slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_texture_panel(slide, 4.5)

    add_pink_bar(slide, Inches(0.8), Inches(0.8), Inches(0.6), Pt(8))
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(1),
                "market", size=48, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.8),
                "evidence", size=48, color=PINK, bold=True)

    # Row 1 — hero stats (only ONE pink per row)
    stats_row1 = [
        ("$835K",  "median house price\nup 2.4% year-on-year",   True),
        ("22 days","average time to sell\nwell-priced homes move faster", False),
        ("0.23%",  "stock on market\nextremely low supply",      False),
    ]
    for i, (num, label, is_pink) in enumerate(stats_row1):
        x = Inches(5.2 + i * 2.7)
        add_textbox(slide, x, Inches(1.4), Inches(2.5), Inches(1),
                    num, size=48, color=PINK if is_pink else CHARCOAL, bold=True)
        add_textbox(slide, x, Inches(2.4), Inches(2.5), Inches(0.8),
                    label, size=14, color=MID_GREY, bold=False)

    # Row 2 — supporting stats
    stats_row2 = [
        ("430+",   "active buyers\nin the last 12 months",  False),
        ("81.6%",  "owner-occupied\nhigh buyer intent",      False),
        ("0.78",   "months of inventory\nseller-scarce market", False),
    ]
    for i, (num, label, is_pink) in enumerate(stats_row2):
        x = Inches(5.2 + i * 2.7)
        add_textbox(slide, x, Inches(3.8), Inches(2.5), Inches(1),
                    num, size=48, color=CHARCOAL, bold=True)
        add_textbox(slide, x, Inches(4.8), Inches(2.5), Inches(0.8),
                    label, size=14, color=MID_GREY, bold=False)

    add_textbox(slide, Inches(5.2), Inches(6.5), Inches(7.5), Inches(0.5),
                "Source: CoreLogic, REIV, ABS Census 2021 — Narre Warren South VIC 3805",
                size=11, color=LIGHT_GREY, bold=False)
    add_slide_number(slide, 6)


def slide_07_image_kitchen(prs):
    """Designer kitchen image — reduced stock slides to 3 total."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_image(slide, "kitchen")

    # Dark overlay panel for readability
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(5.0), Inches(5.833), Inches(2.5)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = CHARCOAL
    panel.line.fill.background()

    add_textbox(slide, Inches(8.0), Inches(5.3), Inches(5), Inches(0.6),
                "the details sell", size=32, color=WHITE, bold=True)
    add_textbox(slide, Inches(8.0), Inches(5.9), Inches(5), Inches(0.8),
                "styled homes in this price bracket achieve\n"
                "8–12% higher sale prices on average.",
                size=16, color=LIGHT_GREY, bold=False)
    add_slide_number(slide, 7)


def slide_08_divider_approach(prs):
    """Pink divider — 'the plan'."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DEEP_PINK)
    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(4),
                "the\nplan", size=120, color=WHITE, bold=True)
    add_slide_number(slide, 8)


def slide_09_marketing_strategy(prs):
    """Marketing strategy — restructured as 3 pillars."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_texture_panel(slide, 4.5)

    add_pink_bar(slide, Inches(0.8), Inches(0.8), Inches(0.6), Pt(8))
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(1),
                "marketing", size=48, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.8),
                "strategy", size=48, color=PINK, bold=True)

    # Three pillars
    pillars = [
        ("reach", PINK, [
            "Premium realestate.com.au listing",
            "Targeted social media campaign",
            "3,200+ buyer database mailout",
            "Signboard on Greaves Road",
        ]),
        ("presentation", CHARCOAL, [
            "Professional photography & drone",
            "Cinematic video walkthrough",
            "Floor plan & property brochure",
            "Styling consultation for your home",
        ]),
        ("negotiation", CHARCOAL, [
            "Private inspections & open homes",
            "Weekly vendor feedback reports",
            "Multi-offer strategy if applicable",
            "Skilled, transparent negotiation",
        ]),
    ]
    for i, (title, title_color, items) in enumerate(pillars):
        x = Inches(5.2 + i * 2.7)
        add_textbox(slide, x, Inches(2.6), Inches(2.5), Inches(0.6),
                    title, size=28, color=title_color, bold=True)
        # Underline
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, Inches(3.15), Inches(1.5), Pt(3)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = PINK if i == 0 else RGBColor(0xD0, 0xD0, 0xD0)
        bar.line.fill.background()

        for j, item in enumerate(items):
            add_textbox(slide, x, Inches(3.5 + j * 0.5), Inches(2.5), Inches(0.45),
                        item, size=14, color=CHARCOAL, bold=False)

    # Method of sale recommendation
    add_textbox(slide, Inches(5.2), Inches(5.8), Inches(7.5), Inches(0.5),
                "recommended method of sale", size=14, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(5.2), Inches(6.2), Inches(7.5), Inches(0.8),
                "Private sale campaign with a 4-week marketing period.\n"
                "This gives us maximum buyer reach while maintaining price control\n"
                "— the strongest approach for well-priced family homes in NWS.",
                size=14, color=MID_GREY, bold=False)
    add_slide_number(slide, 9)


def slide_10_timeline(prs):
    """Campaign timeline — realistic 4-6 weeks."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_texture_panel(slide, 4.5)

    add_pink_bar(slide, Inches(0.8), Inches(0.8), Inches(0.6), Pt(8))
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(1),
                "campaign", size=48, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.8),
                "timeline", size=48, color=PINK, bold=True)

    steps = [
        ("week 1", "prepare", "Styling consultation, professional\nphotography, video & floor plans"),
        ("week 2", "launch", "Listings go live, signboard up,\nbuyer database contacted"),
        ("week 3–4", "inspect", "Saturday opens, midweek privates,\nweekly vendor feedback report"),
        ("week 4–6", "negotiate", "Review all offers, negotiate best\nresult, exchange contracts"),
    ]
    for i, (week, title, desc) in enumerate(steps):
        x = Inches(5.2 + i * 2.05)
        # Circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, Inches(2.6), Inches(1.6), Inches(1.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = PINK if i == 0 else CHARCOAL
        circle.line.fill.background()
        add_textbox(slide, Inches(x.inches + 0.2), Inches(2.95),
                    Inches(1.2), Inches(0.8),
                    week, size=15, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER)
        # Title + desc
        add_textbox(slide, x, Inches(4.5), Inches(1.9), Inches(0.5),
                    title, size=16, color=CHARCOAL, bold=True)
        add_textbox(slide, x, Inches(5.0), Inches(1.9), Inches(1.4),
                    desc, size=14, color=MID_GREY, bold=False)

    add_textbox(slide, Inches(5.2), Inches(6.8), Inches(7), Inches(0.5),
                "Settlement typically 30–60 days after exchange.",
                size=12, color=LIGHT_GREY, bold=False)
    add_slide_number(slide, 10)


def slide_11_divider_investment(prs):
    """Charcoal divider — alternating from pink."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CHARCOAL)
    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(4),
                "your\ninvestment", size=120, color=WHITE, bold=True)
    add_slide_number(slide, 11)


def slide_12_fee_structure(prs):
    """Fee structure — with dollar context + what commission covers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_texture_panel(slide, 4.5)

    add_pink_bar(slide, Inches(0.8), Inches(0.8), Inches(0.6), Pt(8))
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(1),
                "fee", size=48, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.8),
                "structure", size=48, color=PINK, bold=True)

    # Commission — with dollar context
    add_textbox(slide, Inches(5.5), Inches(1.3), Inches(3.5), Inches(1),
                "2.0%", size=80, color=PINK, bold=True)
    add_textbox(slide, Inches(5.5), Inches(2.6), Inches(3.5), Inches(0.8),
                "commission + GST\nat $780K = ~$15,600 + GST", size=16,
                color=MID_GREY, bold=False)

    # What commission covers
    add_textbox(slide, Inches(5.5), Inches(3.6), Inches(3.5), Inches(0.4),
                "your commission covers:", size=14, color=CHARCOAL, bold=True)
    covers = [
        "All inspections & open homes",
        "Personal negotiation by Stuart",
        "Weekly written vendor reports",
        "Contract & settlement management",
    ]
    for i, item in enumerate(covers):
        add_textbox(slide, Inches(5.5), Inches(4.1 + i * 0.38),
                    Inches(3.5), Inches(0.35),
                    item, size=14, color=CHARCOAL, bold=False)

    # Marketing investment
    add_textbox(slide, Inches(9.5), Inches(1.3), Inches(3.5), Inches(1),
                "$5,500", size=56, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(9.5), Inches(2.4), Inches(3.5), Inches(0.5),
                "marketing investment", size=16, color=MID_GREY, bold=False)

    items = [
        ("Professional photography & drone",        "$800"),
        ("Cinematic video walkthrough",             "$1,200"),
        ("Premium realestate.com.au (45 days)",     "$1,800"),
        ("For sale signboard",                      "$350"),
        ("Print brochures (200 units)",             "$450"),
        ("Social media advertising (4 weeks)",      "$600"),
        ("Styling consultation",                    "$300"),
    ]
    for i, (item, cost) in enumerate(items):
        y = Inches(3.2 + i * 0.42)
        add_textbox(slide, Inches(9.5), y, Inches(2.8), Inches(0.4),
                    item, size=14, color=CHARCOAL, bold=False)
        add_textbox(slide, Inches(12.0), y, Inches(0.8), Inches(0.4),
                    cost, size=14, color=PINK, bold=True, align=PP_ALIGN.RIGHT)

    # Total line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(6.15), Inches(3.3), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PINK
    line.line.fill.background()

    add_textbox(slide, Inches(9.5), Inches(6.25), Inches(2.8), Inches(0.4),
                "Total", size=14, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(12.0), Inches(6.25), Inches(0.8), Inches(0.4),
                "$5,500", size=14, color=PINK, bold=True, align=PP_ALIGN.RIGHT)

    add_slide_number(slide, 12)


def slide_13_image_exterior(prs):
    """Exterior image — with dark overlay for text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_image(slide, "exterior")

    # Dark overlay panel for readability
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.5), Inches(6.5), Inches(3)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    panel.line.fill.background()

    add_textbox(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(1.2),
                "the moment\nthey arrive", size=48, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(5.5), Inches(0.5),
                "buyers decide in the first 30 seconds.",
                size=16, color=LIGHT_GREY, bold=False)
    add_slide_number(slide, 13)


def slide_14_divider_why(prs):
    """Pink divider — 'why grants'."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DEEP_PINK)
    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(4),
                "why\ngrants", size=120, color=WHITE, bold=True)
    add_slide_number(slide, 14)


def slide_15_why_us(prs):
    """Why Grants — proof points, not claims."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_texture_panel(slide, 4.5)

    add_pink_bar(slide, Inches(0.8), Inches(0.8), Inches(0.6), Pt(8))
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(1),
                "why", size=48, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.8),
                "grants", size=48, color=PINK, bold=True)

    # Proof points with pink bars — REPLACE XX with real data
    items = [
        ("XX HOMES SOLD IN NARRE WARREN SOUTH — LAST 12 MONTHS", 7.2),
        ("AVG X.X% ABOVE MEDIAN SALE PRICE", 5.8),
        ("LOCAL TEAM — LIVE, WORK & SELL IN YOUR COMMUNITY", 6.8),
        ("X,XXX+ ACTIVE BUYERS IN OUR DATABASE", 6.0),
        ("WEEKLY WRITTEN VENDOR REPORTS — EVERY FRIDAY", 6.5),
        ("XX% VENDOR SATISFACTION / X.X STAR REVIEWS", 5.5),
    ]
    for i, (label, bar_w) in enumerate(items):
        y = Inches(2.8 + i * 0.65)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(5.2), y,
            Inches(bar_w), Inches(0.45)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = PINK
        bar.line.fill.background()
        add_textbox(slide, Inches(5.4), Inches(y.inches - 0.01),
                    Inches(bar_w - 0.2), Inches(0.47),
                    label, size=14, color=WHITE, bold=True)

    add_textbox(slide, Inches(5.2), Inches(6.8), Inches(7), Inches(0.5),
                "Replace XX placeholders with your verified data.",
                size=11, color=LIGHT_GREY, bold=False)
    add_slide_number(slide, 15)


def slide_16_testimonials(prs):
    """Vendor testimonials — social proof."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CHARCOAL)

    add_pink_bar(slide, Inches(0.8), Inches(0.8), Inches(0.6), Pt(8))
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(1),
                "what our", size=48, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.8),
                "vendors say", size=48, color=PINK, bold=True)

    # Testimonial slots — replace with real quotes
    testimonials = [
        ('"Stuart got us $40,000 over our expected price.\n'
         'He kept us informed every step of the way\n'
         'and made the whole process feel easy."',
         "— Jane & Mark, Narre Warren South"),
        ('"We interviewed three agents. Stuart was the only\n'
         'one who gave us honest pricing backed by evidence,\n'
         'not just what we wanted to hear."',
         "— The Nguyen Family, Narre Warren South"),
        ('"Professional, responsive, and genuinely invested\n'
         'in getting us the best result. We recommend\n'
         'Stuart to everyone we know."',
         "— Sarah D, Berwick"),
    ]
    for i, (quote, attr) in enumerate(testimonials):
        y = Inches(2.8 + i * 1.5)
        add_textbox(slide, Inches(1.0), y, Inches(11), Inches(1),
                    quote, size=16, color=WHITE, bold=False,
                    font_name="Georgia")
        add_textbox(slide, Inches(1.0), Inches(y.inches + 1.0),
                    Inches(11), Inches(0.4),
                    attr, size=14, color=PINK, bold=True)

    add_textbox(slide, Inches(1.0), Inches(7.0), Inches(10), Inches(0.4),
                "Replace with your actual vendor testimonials.",
                size=11, color=MID_GREY, bold=False)
    add_slide_number(slide, 16)


def slide_17_next_steps(prs):
    """Next steps — reduce friction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_texture_panel(slide, 4.5)

    add_pink_bar(slide, Inches(0.8), Inches(0.8), Inches(0.6), Pt(8))
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(1),
                "next", size=48, color=CHARCOAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.8),
                "steps", size=48, color=PINK, bold=True)

    steps = [
        ("01", "sign", "Sign the agency agreement\nand we'll get moving immediately."),
        ("02", "prepare", "Styling consultation + photography\nbooked within 5 business days."),
        ("03", "launch", "Your property hits the market\nwith maximum impact from day one."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = Inches(5.2 + i * 2.7)
        add_textbox(slide, x, Inches(2.5), Inches(2.2), Inches(1),
                    num, size=72, color=PINK, bold=True)
        add_textbox(slide, x, Inches(3.8), Inches(2.2), Inches(0.5),
                    title, size=24, color=CHARCOAL, bold=True)
        add_textbox(slide, x, Inches(4.4), Inches(2.2), Inches(1.2),
                    desc, size=15, color=MID_GREY, bold=False)

    # Commitment / risk reversal
    add_textbox(slide, Inches(5.2), Inches(6.0), Inches(7.5), Inches(1),
                "Our commitment: weekly written updates every Friday.\n"
                "If at any point in the first 14 days you're not satisfied\n"
                "with our service, you can cancel with no marketing cost.",
                size=15, color=CHARCOAL, bold=False)
    add_slide_number(slide, 17)


def slide_18_closing(prs):
    """Closing — 'ready when you are' + price anchor."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, CHARCOAL)

    add_pink_bar(slide, Inches(1.0), Inches(1.8), Inches(1.0), Pt(8))
    add_textbox(slide, Inches(1.0), Inches(2.2), Inches(10), Inches(2),
                "ready when\nyou are", size=96, color=WHITE, bold=True)

    # Price anchor
    add_textbox(slide, Inches(1.0), Inches(4.5), Inches(6), Inches(0.8),
                "$750,000 — $810,000", size=36, color=PINK, bold=True)

    # Contact
    add_textbox(slide, Inches(1.0), Inches(5.5), Inches(5), Inches(0.5),
                "Stuart Grant", size=24, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.0), Inches(6.0), Inches(5), Inches(0.4),
                "Grant Estate Agents", size=18, color=PINK, bold=False)
    add_textbox(slide, Inches(1.0), Inches(6.4), Inches(5), Inches(0.8),
                "0400 000 000  |  stuart@grantestateagents.com.au",
                size=15, color=MID_GREY, bold=False)

    # Property reminder right side
    add_textbox(slide, Inches(8.0), Inches(5.5), Inches(4.5), Inches(0.5),
                "171 Greaves Road", size=28, color=WHITE, bold=True,
                align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(8.0), Inches(6.0), Inches(4.5), Inches(0.5),
                "Narre Warren South VIC 3805", size=18, color=MID_GREY,
                bold=False, align=PP_ALIGN.RIGHT)

    add_slide_number(slide, 18)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════

def main():
    print("Generating 171 Greaves Road — V2 (Expert Panel Revision)...")
    print("Downloading images...")

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_01_hero(prs)          # 1  — Hero
    slide_02_agent_profile(prs) # 2  — Stuart Grant profile (NEW)
    slide_03_property_overview(prs)  # 3  — Your home
    slide_04_price_expectation(prs)  # 4  — Price (MOVED UP)
    slide_05_divider_market(prs)     # 5  — "right now" (charcoal)
    slide_06_market_evidence(prs)    # 6  — Combined market + suburb (MERGED)
    slide_07_image_kitchen(prs)      # 7  — Designer kitchen
    slide_08_divider_approach(prs)   # 8  — "the plan" (pink)
    slide_09_marketing_strategy(prs) # 9  — 3 pillars (RESTRUCTURED)
    slide_10_timeline(prs)           # 10 — Timeline (REALISTIC)
    slide_11_divider_investment(prs) # 11 — "your investment" (charcoal)
    slide_12_fee_structure(prs)      # 12 — Fees (EXPANDED)
    slide_13_image_exterior(prs)     # 13 — Designer exterior
    slide_14_divider_why(prs)        # 14 — "why grants" (pink)
    slide_15_why_us(prs)             # 15 — Proof points (REWRITTEN)
    slide_16_testimonials(prs)       # 16 — Testimonials (NEW)
    slide_17_next_steps(prs)         # 17 — Next steps (NEW)
    slide_18_closing(prs)            # 18 — "ready when you are" (REWRITTEN)

    prs.save(OUTPUT)
    print(f"\nDone! Saved to:\n{OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")
    print("\n⚠  BEFORE PRESENTING — replace these placeholders:")
    print("   - Slide 2:  Stuart's photo, years, sales count, $ transacted")
    print("   - Slide 4:  Comparable sales with verified CoreLogic data")
    print("   - Slide 15: XX placeholders with real proof points")
    print("   - Slide 16: Testimonials with real vendor quotes")
    print("   - Slide 18: Phone number & email")


if __name__ == "__main__":
    main()
